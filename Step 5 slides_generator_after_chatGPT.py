from pptx import Presentation

def create_presentation_from_text(file_path, output_path):
    # Create a new PowerPoint presentation
    presentation = Presentation()
    
    # Open and read the text file
    with open(file_path, 'r') as file:
        lines = file.readlines()
    
    slide = None
    for line in lines:
        line = line.strip()
        
        # Detect slide heading
        if line.startswith("### **Slide"):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
        
        # Add title to the slide
        elif line.startswith("**") and line.endswith("**"):
            if slide:
                slide.shapes.title.text = line.strip("**")
        
        # Add content (bullets) to the slide
        elif line.startswith("-"):
            if slide:
                content_box = slide.placeholders[1]
                content_box.text += f"{line[1:].strip()}\n"
    
    # Save the PowerPoint file
    presentation.save(output_path)
    print(f"Presentation saved to {output_path}")

# Example usage
create_presentation_from_text("generated_slides_from_chatGPT.txt", "out_presentation.pptx")
