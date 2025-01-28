import os
from flask import Flask, request, send_file, jsonify, render_template
from pptx import Presentation
import subprocess
import tempfile

# Initialize the Flask application
app = Flask(__name__, template_folder="D:/slides-ai/templates")

# Function to read text from PowerPoint
def read_pptx(file_path):
    presentation = Presentation(file_path)
    text_runs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text.strip())  # Strip unnecessary spaces

    return text_runs

# Function to get summary from the chatbot
def get_summary_from_chatbot(full_text):
    role_instruction = "You are a PowerPoint Slides Improvement agent. You state the heading and the content in each slide as it is relevant to the topic and also add any details that might be informative for the slides."
    full_question = f"{role_instruction} Generate the slides from the following text: {full_text}"
    command = f'ollama run gemma2:2b "{full_question}"'

    process = subprocess.Popen(
        command,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True
    )

    result = []
    try:
        # Read output word by word
        buffer = ""
        while True:
            char = process.stdout.read(1)  # Read one character at a time
            if not char and process.poll() is not None:
                break  # Exit if the process is finished
            if char.isspace():  # Word boundary (space, newline, etc.)
                if buffer:  # If a word has been built up
                    result.append(buffer)
                    buffer = ""  # Reset buffer for the next word
            else:
                buffer += char  # Build up a word
    except UnicodeDecodeError:
        print("A decoding error occurred. Some characters may not be displayed correctly.")
    finally:
        # Ensure the process is closed properly
        process.stdout.close()
        process.stderr.close()
        process.wait()

    return " ".join(result)  # Return the final combined output

# Function to create a new PowerPoint file
def create_ppt(title, content, output_file="output.pptx"):
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Add a slide with a title and content layout
    slide_layout = presentation.slide_layouts[1]  # Use 'Title and Content' layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set the title
    slide.shapes.title.text = title

    # Set the content
    content_box = slide.placeholders[1]
    content_box.text = content

    # Save the PowerPoint file
    presentation.save(output_file)

# Flask route to serve the index.html file
@app.route('/')
def index():
    return render_template('index.html')

# Flask route to upload the PowerPoint file
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    # Save the uploaded file to a temporary location
    temp_dir = tempfile.mkdtemp()
    input_ppt_path = os.path.join(temp_dir, file.filename)
    file.save(input_ppt_path)

    # Step 1: Read text from the input PowerPoint
    text_runs = read_pptx(input_ppt_path)

    # Combine all text into one string
    full_text = "\n".join(text_runs)

    # Step 2: Send the combined text to the chatbot for summarization
    summary = get_summary_from_chatbot(full_text)

    # Step 3: Split the summary into words and remove the first 1 word
    summary_words = summary.split()
    filtered_summary = " ".join(summary_words[1:])  # Remove the first 1 word

    # Step 4: Split the filtered summary into heading and content
    if "*" in filtered_summary:
        heading, content = filtered_summary.split("*", 1)  # Split at the first "*"
        heading = heading.strip()
        content = content.strip()
    else:
        heading = filtered_summary.strip()  # If no "*" is found, take the whole as heading
        content = ""

    # Create output PowerPoint file
    output_ppt_path = os.path.join(temp_dir, "Summarized_Slides.pptx")
    create_ppt(heading, content, output_ppt_path)

    # Return the file for download
    return send_file(output_ppt_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
