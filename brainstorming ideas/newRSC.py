import subprocess
from pptx import Presentation
from ollama import chat

# Function to read text from PowerPoint
def read_pptx(file_path):
    presentation = Presentation(file_path)
    text_runs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text.strip())  # Strip unnecessary spaces

    return text_runs

# Function to get summary from the chatbot
def get_summary_from_chatbot(full_text, output_file="chatbot_summary.txt"):
    role_instruction = "You are a PowerPoint Slides Creator. Generate slides with headings and content for each slide, separated by '*'. Separate each slide with double newlines."
    full_question = f"{role_instruction} Generate case-based learning slides and structure slides from this text: {full_text}. Structure everything by starting with a theoretical case about the subtopic and then provide slides for it and do this for each subtopic."
    command = f'ollama run gemma2:2b "{full_question}"'

    # Stream the chatbot response
    stream = chat(
        model='gemma2:2b',
        messages=[{'role': 'assistant', 'content': full_question}],
        stream=True,
    )

    # Collect streamed chunks into a list
    result = []
    for chunk in stream:
        if 'message' in chunk and 'content' in chunk['message']:
            result.append(chunk['message']['content'])


    with open(output_file, "w", encoding="utf-8") as file:
        file.write(" ".join(result))


    # Combine all chunks into a single string and return
    return " ".join(result)


# Function to create a PowerPoint file with multiple slides
# def create_ppt(slides_data, output_file="output.pptx"):
#     # Create a PowerPoint presentation object
#     presentation = Presentation()

#     # Add slides for each heading-content pair
#     for slide_data in slides_data:
#         heading = slide_data.get("heading", "Untitled Slide")
#         content = slide_data.get("content", "")

#         # Add a slide with a title and content layout
#         slide_layout = presentation.slide_layouts[1]  # Use 'Title and Content' layout
#         slide = presentation.slides.add_slide(slide_layout)

#         # Set the title
#         slide.shapes.title.text = heading

#         # Set the content
#         content_box = slide.placeholders[1]
#         content_box.text = content

#     # Save the PowerPoint file
#     presentation.save(output_file)
#     print(f"\nPowerPoint file saved as '{output_file}'")


def create_ppt_from_text(text, output_file="Generated_Presentation.pptx"):
    """
    Converts structured text into a PowerPoint presentation.

    Args:
    - text (str): The input text containing slide headings and content in a structured format.
    - output_file (str): The file name for the generated PowerPoint presentation.

    Returns:
    - None: Saves the PowerPoint file to the specified location.
    """
    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Split the text into slides based on the delimiter "---"
    slides_data = text.split("---")
    for slide_data in slides_data:
        slide_data = slide_data.strip()
        if not slide_data:
            continue  # Skip empty sections

        # Extract heading and content
        lines = slide_data.split("\n")
        heading = ""
        content = []

        for line in lines:
            line = line.strip()
            if line.startswith("**") and line.endswith("**"):
                heading = line.strip("* ").strip()
            elif line:
                content.append(line.strip("* "))

        # Add a new slide with the heading and content
        if heading or content:  # Ensure there's something to add
            slide_layout = presentation.slide_layouts[1]  # Title and Content layout
            slide = presentation.slides.add_slide(slide_layout)

            # Set the slide title
            slide.shapes.title.text = heading

            # Set the slide content
            content_box = slide.placeholders[1]
            content_box.text = "\n".join(content)

    # Save the PowerPoint presentation
    presentation.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    # File path of the input PowerPoint presentation
    file_path = r"D:\slides-ai\Heart Attack.pptx"  # Replace with your actual file path

    # Step 1: Read text from the input PowerPoint
    text_runs = read_pptx(file_path)

    # Combine all text into one string
    full_text = "\n".join(text_runs)  # Join all slide text into a single string
    print(full_text)

    # Step 2: Send the combined text to the chatbot for slide generation
    print("Processing input to AI for slide generation...")
    ai_response = get_summary_from_chatbot(full_text)
    print(ai_response)


    # Step 4: Create a new PowerPoint file with the parsed slides
    create_ppt_from_text(ai_response, output_file="Generated_Slides.pptx")






