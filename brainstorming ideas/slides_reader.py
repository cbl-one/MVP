from pptx import Presentation

def extract_text_from_pptx(pptx_file, output_txt_file):
    # Open the PowerPoint file
    presentation = Presentation(pptx_file)
    
    with open(output_txt_file, 'w', encoding='utf-8') as output_file:
        for slide_number, slide in enumerate(presentation.slides, start=1):
            output_file.write(f"--- Slide {slide_number} ---\n")
            for shape in slide.shapes:
                # Check if the shape has text
                if hasattr(shape, "text") and shape.text:
                    output_file.write(shape.text + '\n')
            output_file.write('\n')  # Add an empty line after each slide

if __name__ == "__main__":
    # Specify the PowerPoint file and the output text file
    pptx_file = "Heart Attack.pptx"  # Replace with your .pptx file path
    output_txt_file = "output.txt"  # Replace with your desired text file path
    
    extract_text_from_pptx(pptx_file, output_txt_file)
    print(f"Content extracted to {output_txt_file}")
