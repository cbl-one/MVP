import subprocess
from pptx import Presentation

# Function to read text from PowerPoint
def read_pptx(file_path):
    presentation = Presentation(file_path)
    text_runs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text.strip())  # Strip unnecessary spaces

    return text_runs

# Function to get summary from the chatbot
def get_summary_from_chatbot(full_text):
    role_instruction = "You are a PowerPoint Slides Summarizer. You state the heading and the content in max 2 lines"
    full_question = f"{role_instruction} Summarize in max 2 lines - {full_text}"
    command = f'ollama run gemma2:2b "{full_question}"'

    process = subprocess.Popen(
        command,
        shell=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True
    )

    result = []
    try:
        # Read output word by word
        buffer = ""
        while True:
            char = process.stdout.read(1)  # Read one character at a time
            if not char and process.poll() is not None:
                break  # Exit if the process is finished
            if char.isspace():  # Word boundary (space, newline, etc.)
                if buffer:  # If a word has been built up
                    result.append(buffer)
                    buffer = ""  # Reset buffer for the next word
            else:
                buffer += char  # Build up a word
    except UnicodeDecodeError:
        print("A decoding error occurred. Some characters may not be displayed correctly.")
    finally:
        # Ensure the process is closed properly
        process.stdout.close()
        process.stderr.close()
        process.wait()

    return " ".join(result)  # Return the final combined output

# Function to create a new PowerPoint file
def create_ppt(title, content, output_file="output.pptx"):
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Add a slide with a title and content layout
    slide_layout = presentation.slide_layouts[1]  # Use 'Title and Content' layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set the title
    slide.shapes.title.text = title

    # Set the content
    content_box = slide.placeholders[1]
    content_box.text = content

    # Save the PowerPoint file
    presentation.save(output_file)
    print(f"\nPowerPoint file saved as '{output_file}'")

if __name__ == "__main__":
    # File path of the input PowerPoint presentation
    file_path = r"D:\slides-ai\Heart Attack.pptx"  # Replace with your actual file path

    # Step 1: Read text from the input PowerPoint
    text_runs = read_pptx(file_path)

    # Combine all text into one string
    full_text = "\n".join(text_runs)  # Join all slide text into a single string

    # Step 2: Send the combined text to the chatbot for summarization
    print("Processing input to AI for summarization...")
    summary = get_summary_from_chatbot(full_text)

    # Step 3: Split the summary into words and remove the first 22 words
    summary_words = summary.split()
    filtered_summary = " ".join(summary_words[1:])  # Remove the first 22 words
    print(filtered_summary)

    # Step 4: Split the filtered summary into heading and content
    # Everything before the first "*" is the heading and after it is the content
    if "*" in filtered_summary:
        heading, content = filtered_summary.split("*", 1)  # Split at the first "*"
        heading = heading.strip()  # Clean the heading (remove leading/trailing spaces)
        content = content.strip()  # Clean the content
    else:
        heading = filtered_summary.strip()  # If no "*" is found, take the whole as heading
        content = ""

    # Step 5: Create a new PowerPoint file with the cleaned content
    create_ppt(heading, content, output_file="Summarized_Slides.pptx")

