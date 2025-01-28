from pptx import Presentation


def create_ppt_from_text(text, output_file="Generated_Presentation.pptx"):
    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Split the text into slides based on the delimiter "***"
    slides_data = text.split("***")
    for slide_data in slides_data:
        slide_data = slide_data.strip()
        if not slide_data:
            continue  # Skip empty sections

        # Extract heading and content
        lines = slide_data.split("\n")
        heading = ""
        content = []

        for line in lines:
            line = line.strip()
            if line.startswith("**") and line.endswith("**"):
                heading = line.strip("* ").strip()
            elif line:
                content.append(line.strip("* "))

        # Add a new slide with the heading and content
        slide_layout = presentation.slide_layouts[1]  # Title and Content layout
        slide = presentation.slides.add_slide(slide_layout)

        # Set the slide title
        slide.shapes.title.text = heading

        # Set the slide content
        content_box = slide.placeholders[1]
        content_box.text = "\n".join(content)

    # Save the PowerPoint presentation
    presentation.save(output_file)
    print(f"Presentation saved as {output_file}")


def read_text_file(file_path):
    """
    Reads the content of a text file and stores it in a variable.

    Args:
    - file_path (str): The path to the text file.

    Returns:
    - str: The content of the text file as a string.
    """
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            content = file.read()
        return content
    except FileNotFoundError:
        print(f"Error: The file at {file_path} was not found.")
        return ""
    except Exception as e:
        print(f"An error occurred: {e}")
        return ""


def main():
    """
    Main function to read input text and generate a PowerPoint presentation.
    """
    # Path to the input text file
    input_file = "chatbot_summary.txt"  # Replace with your text file path

    # Read the content of the text file
    text_content = read_text_file(input_file)

    if text_content:
        # Generate the PowerPoint presentation
        create_ppt_from_text(text_content, output_file="Generated_Presentation.pptx")
    else:
        print("Failed to generate presentation due to empty or invalid text content.")


if __name__ == "__main__":
    main()
