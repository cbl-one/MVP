from flask import Flask, request, send_file
from pptx import Presentation
import tempfile
import os
from openai import OpenAI

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

# Create upload folder if it doesn't exist
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

# Set your agent endpoint and access key as environment variables in your OS
agent_endpoint = os.getenv("AGENT_ENDPOINT", "https://agent-f1988330ea4a6bd14b63-e2hmw.ondigitalocean.app/api/v1/")
agent_access_key = os.getenv("AGENT_ACCESS_KEY", "ySXLCmm_kIB0YAPWpoWSPtL7qqJWECtZ")

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        if 'file' not in request.files:
            return 'No file uploaded', 400
        file = request.files['file']
        if file.filename.split('.')[-1].lower() != 'pptx':
            return 'Only PPTX files are allowed', 400
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], file.filename))

        try:
            # Extract text from PPTX
            presentation = Presentation(os.path.join(app.config['UPLOAD_FOLDER'], file.filename))
            text = '\n'.join([slide.shapes.title.text for slide in presentation.slides if slide.shapes.title])
        except Exception as e:
            return str(e), 500

        try:
            # Initialize OpenAI client
            client = OpenAI(
                base_url=agent_endpoint,
                api_key=agent_access_key,
            )
            # Send request to endpoint
            response = client.chat.completions.create(
                model="n/a",
                messages=[{"role": "user", "content": text}],
                extra_body={"include_retrieval_info": True}
            )
            # Get the response text
            result_text = response.choices[0].message.content
        except Exception as e:
            return f"Error communicating with endpoint: {str(e)}", 500

        try:
            # Create new presentation with response text
            new_presentation = Presentation()
            # Use the appropriate slide layout for the presentation
            slide_layout = new_presentation.slide_layouts[1]
            
            # Split text into slide contents (using "Slide:" as separator)
            slide_contents = [content.strip() for content in result_text.split("Slide") if content.strip()]
            
            # Create a slide for each content section
            for content in slide_contents:
                slide = new_presentation.slides.add_slide(slide_layout)
                slide.shapes.placeholders[1].text = content

            # Save and send file
            tmp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            new_presentation.save(tmp_file.name)
            return send_file(tmp_file.name, as_attachment=True)
        except Exception as e:
            return f"Error creating presentation: {str(e)}", 500

    return '''
        <!DOCTYPE html>
        <html>
        <body>
            <h1>Upload PPTX File</h1>
            <form action="" method="post" enctype="multipart/form-data">
                <input type="file" name="file" accept=".pptx">
                <input type="submit" value="Upload">
            </form>
        </body>
        </html>
    '''

if __name__ == '__main__':
    app.run(debug=True)